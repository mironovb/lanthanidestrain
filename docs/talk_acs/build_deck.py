#!/usr/bin/env python3
"""Build the ACS oral presentation.

Assertion-evidence format: every slide title is a sentence stating the
finding, the body is one figure, and the depth is in the speaker notes.
16:9. Figures come from docs/talk_acs (talk-sized) and docs/figures_vogiatzis
(chemistry) and docs/figures_arch.

Usage:  PYTHONPATH=$PWD python3 docs/talk_acs/build_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from PIL import Image

HERE = Path(__file__).resolve().parent
REPO = HERE.parents[1]
FIGV = REPO / "docs/figures_vogiatzis"
FIGA = REPO / "docs/figures_arch"
OUT = REPO / "docs/ACS_talk_2026-08.pptx"

NAVY = RGBColor(0x1A, 0x2E, 0x4A)
INK = RGBColor(0x10, 0x12, 0x14)
SUB = RGBColor(0x52, 0x51, 0x4E)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)

SW, SH = Inches(13.333), Inches(7.5)


def add(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=1.0):
    """runs: list of (text, size, bold, color) -> one paragraph each."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, size, bold, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def title_slide(prs, title, subtitle, kicker=None):
    s = add(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    runs = []
    if kicker:
        runs.append((kicker, 15, False, GREEN))
    runs.append((title, 34, True, WHITE))
    runs.append((subtitle, 17, False, RGBColor(0xC9, 0xD4, 0xE2)))
    textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.2), runs,
            spacing=1.25)
    return s


def section_slide(prs, number, title, bullets):
    s = add(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    bar = s.shapes.add_shape(1, 0, Inches(2.6), Inches(0.16), Inches(2.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    runs = [(number, 15, True, GREEN), (title, 30, True, NAVY)]
    for b in bullets:
        runs.append(("— " + b, 16, False, SUB))
    textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.6), runs,
            spacing=1.3)
    return s


def figure_slide(prs, title, fig_path, notes, caption=None, top=Inches(1.30)):
    s = add(prs)
    textbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.85),
            [(title, 23, True, NAVY)], spacing=1.0)
    img = Image.open(fig_path)
    ar = img.size[0] / img.size[1]
    bottom = Inches(0.95) if caption else Inches(0.45)
    max_h = SH - top - bottom
    max_w = SW - Inches(1.1)
    h = max_h
    w = Emu(int(h * ar))
    if w > max_w:
        w = max_w
        h = Emu(int(w / ar))
    left = Emu(int((SW - w) / 2))
    s.shapes.add_picture(str(fig_path), left, top, width=w, height=h)
    if caption:
        textbox(s, Inches(0.55), top + h + Inches(0.12), Inches(12.3),
                Inches(0.6), [(caption, 13, False, SUB)])
    s.notes_slide.notes_text_frame.text = notes
    return s


def table_slide(prs, title, headers, rows, notes, widths=None,
                highlight=None, caption=None):
    s = add(prs)
    textbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.85),
            [(title, 23, True, NAVY)], spacing=1.0)
    n = len(rows) + 1
    tbl = s.shapes.add_table(n, len(headers), Inches(0.75), Inches(1.45),
                             Inches(11.8), Inches(0.42 * n)).table
    if widths:
        total = sum(widths)
        for j, wdt in enumerate(widths):
            tbl.columns[j].width = Emu(int(Inches(11.8) * wdt / total))
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xE3, 0xE5, 0xE8)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        hl = highlight is not None and (i - 1) in highlight
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if hl else WHITE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(13.5)
            r.font.bold = hl
            r.font.color.rgb = INK
    if caption:
        # rows auto-expand for multi-line cells, so pin the caption low
        cap_y = max(Inches(1.45) + Inches(0.52 * n) + Inches(0.3),
                    Inches(5.9))
        textbox(s, Inches(0.75), cap_y, Inches(11.8), Inches(0.9),
                [(caption, 13, False, SUB)])
    s.notes_slide.notes_text_frame.text = notes
    return s


def bullets_slide(prs, title, items, notes, footer=None):
    s = add(prs)
    textbox(s, Inches(0.55), Inches(0.30), Inches(12.3), Inches(0.9),
            [(title, 23, True, NAVY)], spacing=1.0)
    runs = []
    for head, body in items:
        runs.append((head, 18, True, INK))
        runs.append((body, 15, False, SUB))
    textbox(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(5.0), runs,
            spacing=1.25)
    if footer:
        textbox(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6),
                [(footer, 13, False, SUB)])
    s.notes_slide.notes_text_frame.text = notes
    return s


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # 1 ---------------------------------------------------------------
    s = title_slide(
        prs,
        "What limits machine learning of adjacent-lanthanide separation",
        "Bogdan Mironov · with K. Vogiatzis · August 2026",
        kicker="A semiempirical-Hamiltonian bottleneck, and an architecture "
               "that works around it")
    s.notes_slide.notes_text_frame.text = (
        "Two results. First, a chemistry measurement: the semiempirical "
        "method everyone uses for these complexes gets the lanthanide "
        "contraction 2.5x too small, which explains a long run of null 3D "
        "results in the literature. Second, a modelling result: if you split "
        "the prediction into a block level and a within-block shape, 3D "
        "structure finally contributes, and we confirm it on data that took "
        "no part in any choice.")

    # 2 ---------------------------------------------------------------
    figure_slide(
        prs, "Adjacent lanthanides are the industrially hard case",
        FIGV / "g1_pairs.png",
        "Rare-earth separation is a solvent-extraction problem. Neighbouring "
        "lanthanides differ by about 0.013 angstrom in ionic radius, so "
        "their distribution coefficients are nearly identical and hundreds "
        "of stages are needed. Predicting the SEPARATION between neighbours "
        "is the quantity of industrial value, and it is much harder than "
        "predicting the absolute distribution coefficient. Note Pm is absent "
        "from all experimental data, so Nd-Sm is never counted as adjacent.",
        caption="Separation factor between neighbours is what industry pays "
                "for — and what models are worst at.")

    # 3 ---------------------------------------------------------------
    table_slide(
        prs, "Data and evaluation: extractants are held out, and we score "
             "separations",
        ["", "value"],
        [["measurements (log D)", "4,746"],
         ["extractants / lanthanides", "162 / 14"],
         ["3D complexes (Architector + GFN2-xTB)", "956"],
         ["cross-validation", "leave-extractants-out, 5 folds × 3 repeats"],
         ["scored quantity",
          "R² of predicted adjacent-pair log SF (905 pairs)"],
         ["zero on this metric", "predicting the average separation"]],
        "The split is by extractant, so an extractant never appears in both "
        "train and test — this is the regime that matters for screening new "
        "ligands. The metric compares two metals measured under the same "
        "extractant and the same conditions, so everything shared by the two "
        "cancels. R-squared of zero means you did as well as predicting the "
        "average separation; published models sat near +0.14 before this "
        "work.",
        widths=[6, 5],
        caption="Everything reported here is out-of-fold under this protocol.")

    # 4 ---------------------------------------------------------------
    section_slide(prs, "PART I", "Do the 3D structures encode the contraction?",
                  ["A measurement on the Hamiltonian, not on the models",
                   "71 ligands × 15 lanthanides × 2 methods, 2,130 optimisations"])

    # 5 ---------------------------------------------------------------
    figure_slide(
        prs, "We measure the contraction as a per-ligand slope, at fixed "
             "ligand environment",
        FIGV / "g5_slope_demo.png",
        "The design: take one ligand, substitute every lanthanide into the "
        "same coordination environment, optimise, and regress the mean "
        "metal-donor distance against the Shannon ionic radius. The slope is "
        "the compliance c_L: 1.00 means the computed structure follows "
        "experiment exactly. Prior benchmarks report per-complex RMSD against "
        "X-ray or DFT, which cannot see a correlated trend error — bonds can "
        "each be acceptable while the series derivative is badly wrong.",
        caption="Slope against Shannon radii; 1.00 would be exact agreement "
                "with experiment.")

    # 6 ---------------------------------------------------------------
    figure_slide(
        prs, "GFN2-xTB under-responds by 2.5×; g-xTB is within 8 % of "
             "experiment",
        FIGV / "f5_contraction.png",
        "GFN2-xTB gives 0.405 +- 0.145, i.e. it reproduces only 40 % of the "
        "real contraction. g-xTB, which puts the f electrons in the valence, "
        "gives 1.078 +- 0.094. The improvement holds on 71 of 71 ligands, "
        "paired p = 5e-52, and replicates in solvent and on an independent "
        "pilot. GFN2's per-ligand slope is also mostly noise: only 23 % of "
        "its non-linear response is shared across ligands, against 96 % for "
        "g-xTB. This is, to our knowledge, the first lanthanide GEOMETRY "
        "validation of g-xTB — its own preprint benchmarks energies.",
        caption="71 ligands, one protocol, one binary. Improves on 71 of 71; "
                "paired p = 4.9 × 10⁻⁵².")

    # 7 ---------------------------------------------------------------
    figure_slide(
        prs, "The cause is in the parameter file: every lanthanide is a "
             "linear interpolation",
        FIGV / "g4_params.png",
        "Every GFN2 parameter from Ce to Lu is linear interpolation between "
        "two fitted anchors, to file precision. Metal identity therefore "
        "enters a GFN2 geometry as a single linear-in-Z scalar: no f shell, "
        "no crystal field, no gadolinium break. That is a design decision by "
        "the method's authors, documented in their own files — not a bug we "
        "found — but it has a consequence nobody had connected: any 3D "
        "descriptor built on these geometries carries one number about the "
        "metal, and that number is already in the model as the tabular ionic "
        "radius.",
        caption="All 16 parameter families are linear in Z to rounding "
                "precision.")

    # 8 ---------------------------------------------------------------
    bullets_slide(
        prs, "The prediction this makes: 3D encoders should be interchangeable "
             "— and they are",
        [("Eight independent 3D encoders span an effective rank of 1.05 of 8.",
          "Architecture changes move the predictions no more than reseeding "
          "the same architecture does."),
         ("At the pair level, the two best encoders correlate at r = 0.963.",
          "A simplicial network over Vietoris–Rips complexes and a distance "
          "graph network make almost the same predictions."),
         ("Better geometry does not help either.",
          "g-xTB geometries: +0.004, not significant. Structures rebuilt in "
          "exact correspondence (455× cleaner): −0.013.")],
        "This is the retrodiction that makes the chemistry finding "
        "load-bearing rather than a curiosity. If metal identity enters the "
        "geometry as one scalar, then every encoder reading those geometries "
        "is reading the same one scalar, and no architecture can rescue it. "
        "We measured that directly: effective rank 1.05 of 8. And when we "
        "swapped in the physically correct g-xTB geometries, the score did "
        "not move, because 96 % of what changed is a pure function of metal "
        "identity, which the model already had.",
        footer="So the ceiling on the 3D channel is set by the Hamiltonian, "
               "not by the network.")

    # 9 ---------------------------------------------------------------
    section_slide(prs, "PART II", "Then how do we make 3D contribute at all?",
                  ["Change what the model is asked to predict",
                   "Level and shape are different problems"])

    # 10 --------------------------------------------------------------
    figure_slide(
        prs, "87 % of the signal is the block level — and the metric reads "
             "none of it",
        HERE / "T1_scoring.png",
        "Here is one extractant at one condition set, with 14 lanthanides. "
        "The block sits at some level, and within the block the metals tilt. "
        "Decomposed over the whole dataset, 87 % of log D variance is the "
        "level and 13 % is the within-block shape. The separation factor is "
        "a difference inside a block, so the level cancels exactly. A model "
        "trained to minimise error on log D spends nearly all its capacity "
        "on the part of the signal that the metric throws away.")

    # 11 --------------------------------------------------------------
    figure_slide(
        prs, "So predict the level and the shape with separate models",
        HERE / "T2_architecture.png",
        "The level model is an ordinary gradient-boosted tree on log D. The "
        "shape model is the same learner trained on a different target: log D "
        "minus its block mean — so it cannot see the loud between-block "
        "signal at all. The 3D encoder is block-centred the same way. The "
        "prediction is the block anchor plus a weighted sum of the two shape "
        "sources. The anchor is constant inside a block, so it cancels in "
        "every scored comparison: it is there to place the prediction, not "
        "to win points.")

    # 12 --------------------------------------------------------------
    figure_slide(
        prs, "On a real block: both shape sources track the tilt, the 3D one "
             "differently",
        HERE / "T3_block.png",
        "Same block as before, now with the block mean removed. Black is "
        "measured, blue the tabular shape, orange the 3D shape, green the "
        "blend. On the right are the twelve adjacent separations this block "
        "contributes to the metric — note Nd-Sm is absent because Pm does not "
        "exist in the data. You can see both the successes and the failures: "
        "Tb-Dy and Dy-Ho have the wrong sign here.")

    # 13 --------------------------------------------------------------
    figure_slide(
        prs, "One model with this split beats every combination we ever "
             "fitted",
        HERE / "T4_scoreboard.png",
        "The flat champion sits at +0.268. Splitting level from shape takes "
        "the same learner and the same features to +0.318 — that alone beats "
        "the previous best system, a pair-fitted combination of three "
        "different models at +0.313. Adding the 3D shape gives +0.327. Two "
        "things to notice: the architectural effect is larger than the 3D "
        "effect, and the 3D encoder alone is only +0.266, so this is "
        "complementarity, not a strong standalone model.",
        caption="Every bar is out-of-fold, leave-extractants-out, on the same "
                "905 adjacent pairs.")

    # 14 --------------------------------------------------------------
    figure_slide(
        prs, "The 3D weight is fitted, has an interior optimum, and picks the "
             "distance encoder",
        HERE / "T5_weight.png",
        "Panel A: sweeping the mixing weight gives a smooth interior optimum "
        "near 0.3 for the distance encoder, and a monotone decline for the "
        "simplicial one. Panel B: chosen nested per held-out extractant, the "
        "weight lands at 0.35 with a narrow range, so it is not a fragile "
        "choice. Panel C is the mechanism: both encoders correlate with what "
        "the tabular model misses at about 0.2, but they correlate with each "
        "other at 0.963, so the blend takes one and ignores the other. This "
        "is the pair-level version of the rank-1.05 result from Part I.")

    # 15 --------------------------------------------------------------
    figure_slide(
        prs, "The gain holds on 444 pairs that took no part in any decision",
        HERE / "T6_confirm.png",
        "The honesty problem: we had tuned on the same 905 pairs for months. "
        "So we froze a population that had never been used — 444 adjacent "
        "pairs that appear only when the geometry quality filter is relaxed — "
        "and we committed the decision rule to git, with the weight fixed at "
        "0.35, before the deciding model finished training. The contrast is "
        "+0.016 there, the same size as on the legacy set. It also holds on "
        "the collaborator's independently rebuilt dataset, including his 345 "
        "brand-new pairs, and on independent halves of both seed ensembles.",
        caption="Pre-registered rule, fixed weight, one look.")

    # 16 --------------------------------------------------------------
    figure_slide(
        prs, "Being honest about size: the 3D gain is small and unevenly "
             "spread",
        FIGA / "a3_where.png",
        "This is what a +0.008 average actually looks like. Only 481 of 905 "
        "pairs improve, 6 of 12 series positions, and 36 of 76 extractants; "
        "the top five extractants supply half the gross gain. It is a real "
        "effect that survives a pre-registered held-out check, and it is a "
        "small one. Panel A also shows the bigger remaining problem: the "
        "regression slope of predicted on measured is 0.27, so we "
        "systematically under-predict the largest separations — which is "
        "where most of the remaining error sits.")

    # 17 --------------------------------------------------------------
    section_slide(prs, "PART III", "What does not work — four negative results",
                  ["Each one was expected to work",
                   "Each was tested at the same rigour as the positive"])

    # 18 --------------------------------------------------------------
    figure_slide(
        prs, "Topological representations add nothing beyond the distance "
             "encoder",
        HERE / "T7_negatives.png",
        "Given the same slot, the simplicial network gets a fitted weight of "
        "0.01 and contributes nothing — it is 96 % redundant with the "
        "distance encoder. Triangles are not better than edges at matched "
        "seeds. Persistence descriptors are worse than useless: 22 "
        "persistence statistics take the shape model from +0.318 to +0.090, "
        "and the persistence images to +0.002. The bottom bar is the control "
        "that explains it: replace each column by its within-block mean — "
        "same columns, same between-block information, no within-block "
        "variation — and 78 % of the damage disappears. The harm is "
        "specifically in how these descriptors vary across metals inside a "
        "block, which is the only thing the shape model looks at.",
        caption="I expected the split to protect against this. It does the "
                "opposite.")

    # 19 --------------------------------------------------------------
    table_slide(
        prs, "Three more negatives, each with the check that killed it",
        ["what we tried", "result", "why it fails"],
        [["Conditions as within-pair differences",
          "+0.004", "largest raw correlate in the data (|r| ≈ 0.36), but "
                    "extractant-specific — does not transfer out of fold"],
         ["Training directly on all 6,389 within-block pairs",
          "+0.191 alone", "7× more supervision, strong standalone model, but "
                          "its errors duplicate the existing arms"],
         ["Quantum energies: g-xTB metal-swap series",
          "p 0.010 → 0.15", "pilot signal did not survive a 14,000-point "
                            "frozen-cage replication with a matched control"]],
        "Each of these had a good prior. The conditions channel is the "
        "strongest raw correlate anywhere in our data, and it still does not "
        "generalise to unseen extractants. The all-pairs model is the one "
        "whose training population matches the metric exactly, and it learns "
        "nothing the row models did not already know. The energy result is "
        "the one I most wanted to be true: a pilot on relaxed structures gave "
        "p = 0.010, so we ran 14,000 single points overnight on frozen cages "
        "with the other Hamiltonian as control, and the effect attenuated to "
        "p = 0.15. Same sign, no significance.",
        widths=[5, 2.5, 7],
        caption="Every one of these is in the register with the data that "
                "killed it.")

    # 20 --------------------------------------------------------------
    section_slide(prs, "PART IV", "Independent check, and honest limits",
                  ["A collaborator built a different model on the same data",
                   "What the labels can and cannot support"])

    # 21 --------------------------------------------------------------
    figure_slide(
        prs, "An independently built model reproduces — and ties with ours",
        HERE / "T8_collab.png",
        "A collaborator built a completely different system on the same "
        "measurements: an antisymmetric pair-feature random forest with its "
        "own cohort definition and its own metric. Panel A: we reimplemented "
        "his pipeline from his written spec alone — his repository was not "
        "on our machine — and reproduced his headline numbers within his own "
        "stated tolerance. Panel B: scored on his cohort with his metric, his "
        "model and ours are statistically indistinguishable; every paired "
        "bootstrap interval spans zero. Ours does it under a stricter "
        "cross-validation. Two independent pipelines converging on the same "
        "number is evidence about the problem, not about either model.",
        caption="Reproduced from the written spec; the head-to-head "
                "difference is within noise.")

    # 22 --------------------------------------------------------------
    bullets_slide(
        prs, "What the labels can support — and a claim we had to withdraw",
        [("We previously quoted a ceiling of R² ≈ 0.53. It is not defensible.",
          "The estimator behind it divided the noise of one subset by the "
          "variance of a different, wider one; on its own subset it implies a "
          "ceiling below the score models already achieve."),
         ("What is defensible: a separation reproduces to ~0.16 log units.",
          "Measured across independent condition sets, against a spread of "
          "~0.22 — so the noise cancellation on differencing is real, roughly "
          "6× tighter than the levels it is computed from."),
         ("But no point ceiling is identifiable from this dataset.",
          "70 % of scored pairs come from cells with no replicate at all, and "
          "the replicated cells are selected non-randomly. A ceiling needs "
          "designed replication, not opportunistic duplicates.")],
        "I want to flag this because we had it in earlier talks. The 0.53 "
        "number was withdrawn inside our own repository by a later analysis "
        "and kept being quoted anyway. The honest statement is the middle "
        "one: separations are far more reproducible than the levels they "
        "come from, there is clearly headroom above 0.33, and we cannot put "
        "a number on the limit with this data.",
        footer="Correcting our own published-in-talks number, before someone "
               "else has to.")

    # 23 --------------------------------------------------------------
    bullets_slide(
        prs, "Summary",
        [("The semiempirical geometries barely encode the contraction.",
          "GFN2-xTB reproduces 40 % of it; g-xTB is within 8 %. The cause is "
          "linear-in-Z parameterisation, and it explains why 3D encoders are "
          "interchangeable and why better geometry does not help."),
         ("Separating level from shape is worth more than any representation "
          "change.",
          "+0.268 → +0.318 with the same learner and the same features; that "
          "single model beats our previous best combination of three."),
         ("3D contributes only through the shape channel: +0.008, confirmed "
          "on held-out pairs.",
          "First configuration in this project where 3D survives next to the "
          "strongest tabular model instead of being absorbed by it."),
         ("Four representation ideas failed, and the controls say why.",
          "Topology is redundant; persistence descriptors inject within-block "
          "noise; conditions do not transfer; the energy signal did not "
          "replicate.")],
        "If you take one thing away: on this problem, what you ask the model "
        "to predict matters more than what you feed it. And the reason 3D "
        "underperforms is not the network — it is that the electronic "
        "structure method hands you one number about the metal.",
        footer="Next: designed replication for a real ceiling; the "
               "under-prediction of large separations; Ln-xTB as a third "
               "Hamiltonian.")

    # 24 --------------------------------------------------------------
    s = add(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(3.0),
            [("Thank you", 34, True, WHITE),
             ("Data: SAFE solvent-extraction database · Structures: "
              "Architector + GFN2-xTB / g-xTB", 15, False,
              RGBColor(0xC9, 0xD4, 0xE2)),
             ("Code, out-of-fold predictions and every negative result: "
              "github.com/mironovb/lanthanidestrain", 15, False,
              RGBColor(0xC9, 0xD4, 0xE2))], spacing=1.4)
    s.notes_slide.notes_text_frame.text = (
        "Questions I expect: (1) why not DFT — cost, 956 complexes times 15 "
        "metals; and the point is precisely what the cheap method does to "
        "downstream ML. (2) Is +0.008 worth it — on its own, marginal; as "
        "evidence about where 3D can act, it is the whole point. (3) Ln-xTB "
        "(Zhang 2026) is the obvious third Hamiltonian and we have not run "
        "it yet.")

    # --- backup ------------------------------------------------------
    section_slide(prs, "BACKUP", "Supporting material", [])
    figure_slide(
        prs, "Backup: the evaluation protocol and pair construction",
        FIGV / "g2_cv.png",
        "Leave-extractants-out, 5 folds x 3 repeats. Replicates within a "
        "(block, metal) cell are averaged before differencing, so a metal "
        "measured ten times does not dominate one measured once.")
    figure_slide(
        prs, "Backup: what the 3D encoder actually reads",
        FIGV / "g3_vr.png",
        "Vietoris-Rips complex over the heavy atoms of one complex, with "
        "message passing over edges and (for the simplicial variant) "
        "triangles, then a metal-shell readout.")
    figure_slide(
        prs, "Backup: how the mixing weight and the evidence were computed",
        FIGA / "a2_evidence.png",
        "Full evidence panel: weight curve, nested selection, encoder "
        "correlation structure, cross-population gains, seed halves, and the "
        "representations that fail in the same slot.")
    table_slide(
        prs, "Backup: the label-side series shape",
        ["quantity", "value"],
        [["LOEO R² from pair identity alone", "+0.066"],
         ["split-half reliability of the 12-position profile", "r = 0.75"],
         ["structure beyond a radius ramp", "p = 2 × 10⁻¹³"],
         ["only positive mean position", "Eu–Gd"],
         ["largest negative position", "Gd–Tb"]],
        "A 12-value lookup on pair identity alone reaches +0.066 "
        "out-of-fold. The profile is reproducible across extractant halves "
        "and is not a smooth radius ramp — it has structure at the "
        "half-shell. The learned models already contain this, but as a "
        "data observation it may be of independent interest.",
        widths=[8, 3])

    prs.save(OUT)
    return prs


if __name__ == "__main__":
    p = build()
    n = len(p.slides._sldIdLst)
    notes = sum(1 for s in p.slides
                if s.has_notes_slide and s.notes_slide.notes_text_frame.text)
    pics = sum(len([sh for sh in s.shapes if sh.shape_type == 13])
               for s in p.slides)
    print(f"wrote {OUT}")
    print(f"slides {n} · with notes {notes} · figures embedded {pics}")
